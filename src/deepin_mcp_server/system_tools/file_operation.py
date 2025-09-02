import os
import subprocess
import shutil
from pathlib import Path
import logging
import pandas as pd
import docx
import pptx
import PyPDF2

# Configure logging
logger = logging.getLogger(__name__)

def _open_file(file_path: str) -> str:
    try:
        # 检查文件是否存在
        if not os.path.exists(file_path):
            return f"文件不存在: {file_path}"
            
        # 获取用户环境变量
        user_env = os.environ.copy()
        
        # 确保 DISPLAY 环境变量存在
        if 'DISPLAY' not in user_env:
            user_env['DISPLAY'] = ':0'
            
        # 使用用户环境启动 xdg-open，不等待它完成
        subprocess.Popen(
            ['xdg-open', file_path],
            env=user_env,
            # 重定向标准流以避免管道阻塞
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
            # 分离进程，使其成为独立的进程组
            start_new_session=True
        )
        
        return f"已请求打开文件: {file_path}"
            
    except Exception as e:
        return f"打开文件失败: {str(e)}"

def _copy_file(source_path: str, destination_path: str) -> str:
    try:
        shutil.copy2(source_path, destination_path)
        return f"文件已成功复制到: {destination_path}"
    except Exception as e:
        return f"复制文件失败: {str(e)}"

def _move_file(source_path: str, destination_path: str) -> str:
    try:
        shutil.move(source_path, destination_path)
        return f"文件已成功移动到: {destination_path}"
    except Exception as e:
        return f"移动文件失败: {str(e)}"

def _rename_file(old_path: str, new_name: str) -> str:
    try:
        old_file = Path(old_path)
        new_path = old_file.parent / new_name
        old_file.rename(new_path)
        return f"文件已成功重命名为: {new_name}"
    except Exception as e:
        return f"重命名文件失败: {str(e)}"

def _delete_file(file_path: str) -> str:
    try:
        os.remove(file_path)
        return f"文件已成功删除: {file_path}"
    except Exception as e:
        return f"删除文件失败: {str(e)}"

def _create_file(file_path: str, content: str = "") -> str:
    try:
        with open(file_path, 'w', encoding='utf-8') as f:
            f.write(content)
        return f"文件已成功创建: {file_path}"
    except Exception as e:
        return f"创建文件失败: {str(e)}"

def _create_folder(folder_path: str) -> str:
    try:
        os.makedirs(folder_path, exist_ok=True)
        return f"文件夹已成功创建: {folder_path}"
    except Exception as e:
        return f"创建文件夹失败: {str(e)}"
    
def _batch_rename(folder_path: str, new_name: str) -> str:
    try:
        # 检查文件夹是否存在
        if not os.path.exists(folder_path) or not os.path.isdir(folder_path):
            return f"文件夹不存在或不是有效的文件夹: {folder_path}"
            
        # 获取文件夹中的所有文件
        files = [f for f in os.listdir(folder_path) if os.path.isfile(os.path.join(folder_path, f))]
        
        if not files:
            return f"文件夹为空: {folder_path}"
            
        renamed_count = 0
        errors = []
        
        for index, old_file in enumerate(files, 1):
            try:
                # 获取文件扩展名
                _, ext = os.path.splitext(old_file)
                
                # 构建新文件名
                if len(files) == 1:
                    # 如果只有一个文件，直接使用新名称
                    new_file = f"{new_name}{ext}"
                else:
                    # 如果有多个文件，添加数字后缀
                    new_file = f"{new_name}_{index}{ext}"
                
                # 构建完整的文件路径
                old_path = os.path.join(folder_path, old_file)
                new_path = os.path.join(folder_path, new_file)
                
                # 重命名文件
                os.rename(old_path, new_path)
                renamed_count += 1
                
            except Exception as e:
                errors.append(f"重命名文件 {old_file} 失败: {str(e)}")
        
        # 构建返回消息
        if errors:
            return f"批量重命名完成，成功重命名 {renamed_count} 个文件，失败 {len(errors)} 个文件。\n错误详情：\n" + "\n".join(errors)
        else:
            return f"批量重命名完成，成功重命名 {renamed_count} 个文件"
            
    except Exception as e:
        return f"批量重命名失败: {str(e)}"
    
def _list_dir(folder_path: str, recursive: bool = False) -> str:
    def _list_dir_recursive(path: str) -> tuple[list[str], list[str]]:
        """递归列出所有文件和目录"""
        dirs = []
        files = []
        
        for entry in os.listdir(path):
            if entry.startswith('.'):
                continue
                
            full_path = os.path.join(path, entry)
            if os.path.isdir(full_path):
                dirs.append(entry)
                # 递归处理子目录
                sub_dirs, sub_files = _list_dir_recursive(full_path)
                dirs.extend(f"{entry}/{sub_dir}" for sub_dir in sub_dirs)
                files.extend(f"{entry}/{sub_file}" for sub_file in sub_files)
            else:
                files.append(entry)
                
        return dirs, files

    try:
        # 检查文件夹是否存在
        if not os.path.exists(folder_path) or not os.path.isdir(folder_path):
            return f"文件夹不存在或不是有效的文件夹: {folder_path}"
            
        if recursive:
            dirs, files = _list_dir_recursive(folder_path)
        else:
            # 获取文件夹中的所有条目，过滤隐藏文件
            entries = [e for e in os.listdir(folder_path) if not e.startswith('.')]
            
            if not entries:
                return f"文件夹为空: {folder_path}"
                
            # 分离文件夹和文件
            dirs = []
            files = []
            
            for entry in entries:
                full_path = os.path.join(folder_path, entry)
                if os.path.isdir(full_path):
                    dirs.append(entry)
                else:
                    files.append(entry)
        
        # 组合结果
        result = ""
        if dirs:
            result = "\n".join(dirs)
            result = f"<dir>{result}</dir>"
        if files:
            file = "\n".join(files)
            result = result + f"<file>{file}</file>"
            
        return result if result else f"文件夹为空: {folder_path}"

    except Exception as e:
        return f"获取文件列表失败: {str(e)}"
    
def _read_document(document_path: str) -> str:
    try:
        # 检查文件是否存在
        if not os.path.exists(document_path):
            return f"文件不存在: {document_path}"
            
        # 获取文件扩展名
        _, ext = os.path.splitext(document_path)
        ext = ext.lower().strip('.')
        
        # 检查文件大小
        file_size = os.path.getsize(document_path)
        max_size = 10 * 1024 * 1024  # 10 MB
        if file_size > max_size:
            return f"文件太大，不支持读取超过10MB的文件: {document_path}"
            
        # 处理纯文本文件
        if ext in ['txt', 'md', 'json', 'csv', 'log', 'ini', 'conf']:
            try:
                with open(document_path, 'r', encoding='utf-8') as f:
                    content = f.read()
                return f"<document_content>{content}</document_content>"
            except UnicodeDecodeError:
                try:
                    # 尝试使用其他编码
                    with open(document_path, 'r', encoding='gbk') as f:
                        content = f.read()
                    return f"<document_content>{content}</document_content>"
                except UnicodeDecodeError:
                    return "文件解码失败，可能不是文本文件或者使用了不支持的编码"
        
        # 处理 PDF 文件
        if ext == 'pdf':
            try:
                # 动态导入 PyPDF2 库
                try:
                    import PyPDF2
                except ImportError:
                    return "请安装 PyPDF2 库以读取 PDF 文件：pip install PyPDF2"
                
                text = []
                with open(document_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    for page_num in range(len(pdf_reader.pages)):
                        page = pdf_reader.pages[page_num]
                        text.append(page.extract_text())
                
                return f"<document_content>{''.join(text)}</document_content>"
            except Exception as e:
                logger.error(f"读取PDF文件失败: {e}")
                return f"读取PDF文件失败: {str(e)}"
        
        # 处理 Word 文档 (docx)
        if ext in ['docx', 'doc']:
            try:
                # 动态导入 python-docx
                try:
                    import docx
                except ImportError:
                    return "请安装 python-docx 库以读取 DOCX 文件：pip install python-docx"
                
                # 读取 DOCX 文件
                doc = docx.Document(document_path)
                text = [para.text for para in doc.paragraphs]
                return f"<document_content>{' '.join(text)}</document_content>"
            except Exception as e:
                logger.error(f"读取DOCX文件失败: {e}")
                return f"读取DOCX文件失败: {str(e)}"
        
        # 处理 Excel 文件 (xls, xlsx)
        if ext in ['xls', 'xlsx']:
            try:
                # 动态导入 pandas 和相关依赖
                try:
                    import pandas as pd
                except ImportError:
                    return "请安装 pandas 库以读取 Excel 文件：pip install pandas openpyxl xlrd"
                
                # 读取 Excel 文件
                if ext == 'xlsx':
                    df = pd.read_excel(document_path, engine='openpyxl')
                else:
                    df = pd.read_excel(document_path, engine='xlrd')
                
                # 将数据转换为字符串
                content = df.to_string()
                return f"<document_content>{content}</document_content>"
            except Exception as e:
                logger.error(f"读取Excel文件失败: {e}")
                return f"读取Excel文件失败: {str(e)}"
        
        # 处理 PowerPoint 文件 (pptx)
        if ext == 'pptx':
            try:
                # 动态导入 python-pptx
                try:
                    import pptx
                except ImportError:
                    return "请安装 python-pptx 库以读取 PPTX 文件：pip install python-pptx"
                
                # 读取 PPTX 文件
                presentation = pptx.Presentation(document_path)
                text = []
                
                # 提取所有幻灯片中的文本
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text.append(shape.text)
                
                return f"<document_content>{' '.join(text)}</document_content>"
            except Exception as e:
                logger.error(f"读取PPTX文件失败: {e}")
                return f"读取PPTX文件失败: {str(e)}"
        
        # 处理 PowerPoint 文件 (ppt)
        if ext == 'ppt':
            try:
                # 尝试使用 python-ppt 库（如果有的话）
                try:
                    # 对于 Linux 平台，目前没有良好的纯 Python 库处理 PPT 文件
                    # 可以考虑使用 Apache POI 的 Python 绑定，但这超出了基本需求
                    
                    # 提供替代解决方案
                    return "暂不支持直接读取 PPT 文件，建议将其转换为 PPTX 格式，或使用 LibreOffice 转换为 PDF"
                except ImportError:
                    return "不支持直接读取 PPT 文件，建议将其转换为 PPTX 格式"
            except Exception as e:
                logger.error(f"处理PPT文件失败: {e}")
                return f"处理PPT文件失败: {str(e)}"
        
        # 不支持的文件格式
        return f"不支持的文件格式: {ext}"
            
    except Exception as e:
        logger.error(f"读取文档时发生未知错误: {e}")
        return f"读取文档失败: {str(e)}"

def _get_files_size(file_paths: list[str]) -> str:
    try:
        file_sizes = []
        for file_path in file_paths:
            file_size = os.path.getsize(file_path)
            file_sizes.append(f"{file_path}: {file_size} bytes")
        return "\n".join(file_sizes)
    except Exception as e:
        return f"获取文件大小失败: {str(e)}"